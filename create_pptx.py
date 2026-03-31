from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x29, 0xB6, 0xF6)
ACCENT_GREEN = RGBColor(0x4E, 0xCD, 0xC4)
ACCENT_RED = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT_YELLOW = RGBColor(0xE9, 0xC4, 0x6A)
ACCENT_ORANGE = RGBColor(0xF4, 0xA2, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
MED_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x2A, 0x2A, 0x45)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, corner_radius=None):
    if corner_radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = corner_radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.level = 0
    return txBox


# =========================================================================
# SLIDE 1: Title
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "Super Mario SPCS Telemetry", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(1),
             "Real-time Game Telemetry on Snowpark Container Services", font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
             "OpenTelemetry  |  SPCS Event Table  |  Streamlit in Snowsight", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_shape_rect(slide, Inches(4.5), Inches(5.5), Inches(4.3), Inches(0.06), ACCENT_GREEN)

add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.6),
             "Snowflake Solutions Engineering Demo", font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 2: What We Built
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "What We Built", font_size=36, color=WHITE, bold=True)

cards = [
    ("Game Container", "Super Mario Bros HTML5 game\nrunning on SPCS with public\ningress endpoint", ACCENT_RED),
    ("Telemetry Pipeline", "Browser JS hooks capture events,\nPython sidecar exports via\nOpenTelemetry OTLP gRPC", ACCENT_BLUE),
    ("Event Table", "All spans, metrics, and logs\nland in Snowflake's native\nevent table (90+ metric types)", ACCENT_GREEN),
    ("Analytics Dashboard", "Streamlit in Snowsight with\n5 tabs, KPIs, charts, and\nplatform monitoring", ACCENT_YELLOW),
]

for i, (title, desc, accent) in enumerate(cards):
    x = Inches(0.6 + i * 3.15)
    y = Inches(1.6)
    card = add_shape_rect(slide, x, y, Inches(2.9), Inches(3.5), CARD_BG, corner_radius=0.05)
    add_shape_rect(slide, x, y, Inches(2.9), Inches(0.08), accent)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.35), Inches(2.4), Inches(0.5), title, font_size=20, color=accent, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(1.0), Inches(2.4), Inches(2.2), desc, font_size=15, color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(1.5),
             "End-to-end: from a player pressing keys in the browser to real-time charts in Snowsight,\nall running inside Snowflake with zero external infrastructure.",
             font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 3: Architecture
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Architecture Overview", font_size=36, color=WHITE, bold=True)

boxes = [
    (0.5, 2.5, 2.2, 1.2, "Browser\n(telemetry.js)", ACCENT_BLUE),
    (3.3, 2.5, 2.2, 1.2, "nginx\n(port 8080)", ACCENT_GREEN),
    (6.1, 1.8, 2.2, 1.2, "Tomcat Game\n(port 8888)", ACCENT_RED),
    (6.1, 3.4, 2.5, 1.2, "Python Sidecar\nOpenTelemetry SDK\n(port 9090)", ACCENT_YELLOW),
    (9.3, 3.4, 2.5, 1.2, "SPCS Event Table\nevent_db.event_sh\n.my_events", ACCENT_GREEN),
    (9.3, 1.8, 2.5, 1.2, "Streamlit\nDashboard\n(Snowsight)", ACCENT_BLUE),
]

for (x, y, w, h, text, accent) in boxes:
    card = add_shape_rect(slide, Inches(x), Inches(y), Inches(w), Inches(h), CARD_BG, corner_radius=0.04)
    card.line.color.rgb = accent
    card.line.width = Pt(2)
    tb = add_text_box(slide, Inches(x + 0.1), Inches(y + 0.15), Inches(w - 0.2), Inches(h - 0.2), text, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)

arrows_text = [
    "GET /telemetry?d=...",
    "proxy_pass",
    "Game UI",
    "OTLP gRPC",
    "SQL views",
]
arrow_positions = [
    (1.8, 2.2, 3.5), (4.6, 2.2, 3.5), (7.4, 2.0, 2.0), (7.8, 4.0, 3.8), (10.5, 3.0, 2.0),
]
for txt, (x, y, _) in zip(arrows_text, arrow_positions):
    add_text_box(slide, Inches(x), Inches(y - 0.4), Inches(2), Inches(0.35), txt, font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(1.5),
             "Data Flow: Player actions -> JS hooks -> Tracking pixel GET -> nginx -> Python sidecar -> OpenTelemetry OTLP -> Event Table -> Views -> Streamlit",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 4: Container Stack
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Container Stack (Single SPCS Container)", font_size=36, color=WHITE, bold=True)

layers = [
    ("nginx Reverse Proxy", "Port 8080 (public ingress)", ACCENT_GREEN, 1.5),
    ("Tomcat 9 + Super Mario HTML5", "Port 8888 (game server)", ACCENT_RED, 2.6),
    ("Python Telemetry Sidecar", "Port 9090 (OTel collector)", ACCENT_YELLOW, 3.7),
    ("OpenTelemetry SDK", "OTLP gRPC exporter to SPCS event table", ACCENT_BLUE, 4.8),
]

for (title, desc, accent, y) in layers:
    add_shape_rect(slide, Inches(1.5), Inches(y), Inches(10), Inches(0.9), CARD_BG, corner_radius=0.03)
    add_shape_rect(slide, Inches(1.5), Inches(y), Inches(0.12), Inches(0.9), accent)
    add_text_box(slide, Inches(2.0), Inches(y + 0.1), Inches(4), Inches(0.4), title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(7), Inches(y + 0.15), Inches(4), Inches(0.4), desc, font_size=15, color=MED_GRAY)

add_text_box(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.8),
             "Resources: 500m-1 CPU  |  512Mi-2Gi RAM  |  Compute Pool: CPU_X64_XS (1 node)",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 5: Telemetry Events
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Game Telemetry Events", font_size=36, color=WHITE, bold=True)

events = [
    ("game_start", "Player starts game", "lives"),
    ("level_start", "Level begins", "level, difficulty, type"),
    ("death", "Mario dies", "level, lives, large, fire"),
    ("level_win", "Flag reached", "level, time_left"),
    ("game_over", "All lives lost", "level, coins, duration"),
    ("game_win", "Final level done", "session_duration"),
    ("coin", "Coin collected", "total_coins"),
    ("powerup_spawn", "Block hit", "type, level"),
    ("key_press", "Key pressed", "key (throttled 500ms)"),
]

headers = ["Event", "When", "Key Attributes"]
col_widths = [2.5, 2.8, 4.5]
x_start = 1.5
y_start = 1.4

for j, (hdr, w) in enumerate(zip(headers, col_widths)):
    x = Inches(x_start + sum(col_widths[:j]))
    add_shape_rect(slide, x, Inches(y_start), Inches(w), Inches(0.5), ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.15), Inches(y_start + 0.05), Inches(w - 0.3), Inches(0.4), hdr, font_size=14, color=WHITE, bold=True)

for i, (evt, when, attrs) in enumerate(events):
    y = Inches(y_start + 0.55 + i * 0.48)
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x22, 0x22, 0x3A)
    for j, (val, w) in enumerate(zip([f"mario.{evt}", when, attrs], col_widths)):
        x = Inches(x_start + sum(col_widths[:j]))
        add_shape_rect(slide, x, y, Inches(w), Inches(0.45), bg)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(w - 0.3), Inches(0.35), val, font_size=13, color=LIGHT_GRAY)

# =========================================================================
# SLIDE 6: Platform Metrics
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Platform Metrics (SPCS Native + Derived)", font_size=36, color=WHITE, bold=True)

native = [
    ("container.cpu.usage", "CPU cores", "Container CPU utilization"),
    ("container.memory.usage", "bytes", "Container memory usage"),
    ("storage.used / free", "bytes", "Disk storage utilization"),
    ("system.network.io", "bytes", "Network I/O (rx/tx direction)"),
    ("network.ingress.*", "count", "Active connections & CPS"),
]

derived = [
    ("Disk IOPS", "storage delta / elapsed / 4K", "Estimated I/O operations/sec"),
    ("I/O Throughput", "network delta / elapsed / 1MB", "Estimated throughput MB/s"),
    ("Avg Latency", "elapsed * 1000 / samples", "Estimated I/O latency ms"),
]

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5), "Native SPCS Metrics", font_size=20, color=ACCENT_GREEN, bold=True)

for i, (name, unit, desc) in enumerate(native):
    y = Inches(1.9 + i * 0.55)
    add_shape_rect(slide, Inches(0.8), y, Inches(5.5), Inches(0.5), CARD_BG, corner_radius=0.02)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(2.5), Inches(0.4), name, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(3.5), y + Inches(0.05), Inches(2.5), Inches(0.4), f"{unit} - {desc}", font_size=12, color=MED_GRAY)

add_text_box(slide, Inches(7), Inches(1.3), Inches(5), Inches(0.5), "Derived Metrics (Dashboard)", font_size=20, color=ACCENT_ORANGE, bold=True)

for i, (name, formula, desc) in enumerate(derived):
    y = Inches(1.9 + i * 0.7)
    add_shape_rect(slide, Inches(7), y, Inches(5.5), Inches(0.6), CARD_BG, corner_radius=0.02)
    add_text_box(slide, Inches(7.2), y + Inches(0.02), Inches(2), Inches(0.3), name, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.2), y + Inches(0.3), Inches(5), Inches(0.3), f"{formula}  ->  {desc}", font_size=12, color=MED_GRAY)

add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5),
             "SPCS exposes 90+ distinct metric types via platformMonitor config.\nNo native IOPS/latency counters — derived from storage and network deltas.\nDashboard includes time range selector: 15m / 1h / 6h / 24h / 7d / all time.",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 7: Dashboard Overview
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Streamlit Dashboard", font_size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.5), "KPI Rows", font_size=20, color=ACCENT_BLUE, bold=True)

kpi_row1 = ["Games", "Deaths", "Coins", "Attempts", "Wins", "Powerups", "Keys"]
kpi_row2 = ["Avg Time", "Sec/Death", "Coins/Rnd", "Lvl Time", "Attempts", "Win Rate", "Keys/Game"]

for i, kpi in enumerate(kpi_row1):
    x = Inches(0.8 + i * 1.65)
    add_shape_rect(slide, x, Inches(1.7), Inches(1.5), Inches(0.6), CARD_BG, corner_radius=0.03)
    add_text_box(slide, x + Inches(0.1), Inches(1.78), Inches(1.3), Inches(0.4), kpi, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)

for i, kpi in enumerate(kpi_row2):
    x = Inches(0.8 + i * 1.65)
    add_shape_rect(slide, x, Inches(2.4), Inches(1.5), Inches(0.6), CARD_BG, corner_radius=0.03)
    add_text_box(slide, x + Inches(0.1), Inches(2.48), Inches(1.3), Inches(0.4), kpi, font_size=13, color=ACCENT_YELLOW, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(3.3), Inches(5), Inches(0.5), "5 Dashboard Tabs", font_size=20, color=ACCENT_BLUE, bold=True)

tabs = [
    ("Leaderboard", "Top scores + game sessions\n(ranked by coins, latest first)", ACCENT_RED),
    ("Event Timeline", "Stacked area chart of all\ngame events over time", ACCENT_BLUE),
    ("Deaths & Levels", "Deaths by level + powerups\ncollected breakdown", ACCENT_GREEN),
    ("Controls", "Key press distribution\nbar chart", ACCENT_YELLOW),
    ("Platform Metrics", "CPU, memory, disk, IOPS,\nnetwork with time picker", ACCENT_ORANGE),
]

for i, (title, desc, accent) in enumerate(tabs):
    x = Inches(0.5 + i * 2.55)
    y = Inches(3.9)
    card = add_shape_rect(slide, x, y, Inches(2.4), Inches(2.0), CARD_BG, corner_radius=0.04)
    add_shape_rect(slide, x, y, Inches(2.4), Inches(0.06), accent)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.1), Inches(0.4), title, font_size=16, color=accent, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.7), Inches(2.1), Inches(1.2), desc, font_size=13, color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8),
             "Deployed as Streamlit in Snowsight (SIS) with container runtime + SYSTEM_COMPUTE_POOL_CPU",
             font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 8: Setup Steps
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Setup Steps (Start to Finish)", font_size=36, color=WHITE, bold=True)

steps = [
    ("1", "Infrastructure", "Run 01_infrastructure.sql\nDatabase, warehouse, compute pool,\nevent table, image repo", ACCENT_BLUE),
    ("2", "Docker Build", "Build & push container image\nto Snowflake image registry\n(linux/amd64)", ACCENT_GREEN),
    ("3", "SPCS Service", "Run 02_service.sql\nCreate service with platform\nmonitoring enabled", ACCENT_RED),
    ("4", "Analytics Views", "Run 03_analytics_views.sql\n7 views over the event table\nfor dashboard queries", ACCENT_YELLOW),
    ("5", "Streamlit Deploy", "Run snow streamlit deploy\nor 04_streamlit_dashboard.sql\nDashboard goes live", ACCENT_ORANGE),
    ("6", "Play & Monitor", "Open game URL, play Mario\nWatch real-time telemetry\nin Snowsight dashboard", RGBColor(0x96, 0xCE, 0xB4)),
]

for i, (num, title, desc, accent) in enumerate(steps):
    x = Inches(0.5 + i * 2.1)
    y = Inches(1.5)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), y, Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(24)
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].font.color.rgb = DARK_BG
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, x, Inches(2.4), Inches(2.0), Inches(0.5), title, font_size=17, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.9), Inches(2.0), Inches(2.0), desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
             "All SQL scripts are numbered and idempotent (CREATE OR REPLACE / IF NOT EXISTS).\nTotal setup time: ~15 minutes (excluding Docker build/push).",
             font_size=15, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 9: Key Snowflake Features
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Snowflake Features Demonstrated", font_size=36, color=WHITE, bold=True)

features = [
    ("Snowpark Container Services", "Run any Docker container with public endpoints, auto-scaling, and built-in monitoring"),
    ("SPCS Platform Monitor", "Native CPU, memory, storage, and network metrics with zero configuration"),
    ("Event Tables", "Unified telemetry store for spans, metrics, and logs (OpenTelemetry compatible)"),
    ("Streamlit in Snowsight", "Interactive dashboards with container runtime, compute pools, and pip packages"),
    ("Semi-structured Data", "JSON parsing with record:name::STRING, record_attributes:key::INT syntax"),
    ("Views over Event Tables", "Real-time analytics with SQL views, CTEs, window functions, and pivots"),
]

for i, (title, desc) in enumerate(features):
    y = Inches(1.4 + i * 0.9)
    add_shape_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.75), CARD_BG, corner_radius=0.03)
    add_shape_rect(slide, Inches(0.8), y, Inches(0.1), Inches(0.75), ACCENT_BLUE)
    add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(3.5), Inches(0.35), title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.38), Inches(10.5), Inches(0.35), desc, font_size=14, color=MED_GRAY)

# =========================================================================
# SLIDE 10: Resources
# =========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "Resources & Links", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

resources = [
    "SQL Scripts:  sql/01_infrastructure.sql  ->  04_streamlit_dashboard.sql",
    "Docker Files:  mario-spcs/  (Dockerfile, nginx.conf, telemetry_sidecar.py, telemetry.js)",
    "Dashboard:     mario-streamlit/  (streamlit_app.py, pyproject.toml, snowflake.yml)",
    "README:        Full start-to-finish setup guide with troubleshooting",
    "",
    "Game URL:      https://<ingress-hash>.snowflakecomputing.app",
    "Dashboard:     Snowsight > Streamlit Apps > MARIO_TELEMETRY_DASHBOARD",
]

for i, line in enumerate(resources):
    if line:
        add_text_box(slide, Inches(1.5), Inches(2.8 + i * 0.5), Inches(10), Inches(0.45), line, font_size=16, color=LIGHT_GRAY)

add_shape_rect(slide, Inches(4.5), Inches(6.3), Inches(4.3), Inches(0.06), ACCENT_GREEN)
add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "Built with Cortex Code", font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# =========================================================================
# Save
# =========================================================================
prs.save("/tmp/mario-spcs-docs/Mario_SPCS_Telemetry_Overview.pptx")
print("PowerPoint saved!")
