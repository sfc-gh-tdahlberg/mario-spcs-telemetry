"""
Super Mario SPCS Telemetry — Event Presentation
Uses the Snowflake presentation template (blue backgrounds, white text, Snowflake accent palette).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

TEMPLATE = os.path.join(os.path.dirname(__file__), "presentation_template/presetation_template.pptx")
OUTPUT   = os.path.join(os.path.dirname(__file__), "Super_Mario_SPCS_Telemetry_Presentation.pptx")

prs = Presentation(TEMPLATE)
W = prs.slide_width    # 10 inches
H = prs.slide_height   # 5.62 inches

# ── Remove existing template slides ──────────────────────────────────────────
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
sl_list = prs.slides._sldIdLst
while len(sl_list):
    rId = sl_list[0].get(f'{{{NS_R}}}id')
    prs.part.drop_rel(rId)
    del sl_list[0]

# ── Snowflake Brand Palette (from theme1.xml) ─────────────────────────────────
SF_BLUE   = RGBColor(0x29, 0xB5, 0xE8)   # accent1  — Snowflake primary blue
SF_NAVY   = RGBColor(0x11, 0x56, 0x7F)   # accent2  — Dark navy
SF_TEAL   = RGBColor(0x71, 0xD3, 0xDC)   # accent3  — Teal
SF_ORANGE = RGBColor(0xFF, 0x9F, 0x36)   # accent4  — Orange
SF_PURPLE = RGBColor(0x7D, 0x44, 0xCF)   # accent5  — Purple
SF_PINK   = RGBColor(0xD4, 0x5B, 0x90)   # accent6  — Pink
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xE8, 0xF4, 0xFB)   # very light blue-white
DARK_TEXT = RGBColor(0x26, 0x26, 0x26)   # dk1 from theme

# Background is template blue — all cards/shapes sit on top
CARD_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # white card
CARD_NAVY   = RGBColor(0x11, 0x56, 0x7F)   # navy card
CARD_DEEP   = RGBColor(0x0D, 0x3A, 0x5C)   # deeper navy for contrast
SUBTLE_BG   = RGBColor(0x1A, 0x68, 0x9A)   # slightly lighter navy

# ── Layout shortcuts ──────────────────────────────────────────────────────────
# Layout 0  = 1_Cover 02          (cover, skiing background)
# Layout 6  = 1_Cover 02 6        (content, geometric background)
# Layout 12 = Summit22-Divider    (divider, cloud sky background - great for demo!)
# Layout 13 = Agenda
# Layout 25 = Thank You
LAY_COVER   = prs.slide_layouts[0]
LAY_CONTENT = prs.slide_layouts[6]
LAY_DIVIDER = prs.slide_layouts[12]
LAY_AGENDA  = prs.slide_layouts[13]
LAY_THANKS  = prs.slide_layouts[25]

# ── Primitive helpers ─────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=CARD_NAVY, line_color=None, line_w=1.5, radius=None):
    if radius:
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        sh.adjustments[0] = radius
    else:
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color; sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    return sh

def txt(slide, l, t, w, h, text, sz=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    r = p.runs[0]
    r.font.size = Pt(sz); r.font.color.rgb = color
    r.font.bold = bold;   r.font.italic = italic
    return tb

def multiline_txt(slide, l, t, w, h, lines, sz=13, color=OFF_WHITE, spacing=4):
    """Add multiple lines as separate paragraphs in one text box."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = color
        p.space_after = Pt(spacing)
    return tb

def slide_title(slide, title, subtitle=None, title_y=Inches(0.28), title_sz=28):
    txt(slide, Inches(0.5), title_y, W - Inches(1.0), Inches(0.6),
        title, sz=title_sz, color=WHITE, bold=True)
    if subtitle:
        txt(slide, Inches(0.5), title_y + Inches(0.62), W - Inches(1.0), Inches(0.35),
            subtitle, sz=13, color=OFF_WHITE)

def white_card(slide, l, t, w, h, title, body_lines, accent=SF_TEAL,
               title_sz=13, body_sz=11, radius=0.05):
    """White card with colored top border — Snowflake style."""
    rect(slide, l, t, w, h, fill=CARD_WHITE, radius=radius)
    rect(slide, l, t, w, Inches(0.06), fill=accent, radius=radius)
    txt(slide, l + Inches(0.15), t + Inches(0.12), w - Inches(0.3), Inches(0.38),
        title, sz=title_sz, color=SF_NAVY, bold=True)
    multiline_txt(slide, l + Inches(0.15), t + Inches(0.55), w - Inches(0.3),
                  h - Inches(0.62), body_lines, sz=body_sz, color=DARK_TEXT, spacing=3)

def navy_card(slide, l, t, w, h, title, body_lines, accent=SF_TEAL,
              title_sz=13, body_sz=11, radius=0.05):
    """Navy card — for dark-on-dark layouts."""
    rect(slide, l, t, w, h, fill=CARD_NAVY, radius=radius)
    rect(slide, l, t, Inches(0.08), h, fill=accent, radius=0.02)
    txt(slide, l + Inches(0.2), t + Inches(0.1), w - Inches(0.3), Inches(0.36),
        title, sz=title_sz, color=accent, bold=True)
    multiline_txt(slide, l + Inches(0.2), t + Inches(0.5), w - Inches(0.32),
                  h - Inches(0.58), body_lines, sz=body_sz, color=OFF_WHITE, spacing=3)

def row_card(slide, l, t, w, h, title, desc, accent=SF_TEAL, title_sz=14, desc_sz=11):
    """Horizontal row card (title left, description right)."""
    rect(slide, l, t, w, h, fill=CARD_WHITE, radius=0.04)
    rect(slide, l, t, Inches(0.08), h, fill=accent)
    txt(slide, l + Inches(0.2), t + Inches(0.1), Inches(2.8), h - Inches(0.2),
        title, sz=title_sz, color=SF_NAVY, bold=True)
    txt(slide, l + Inches(3.2), t + Inches(0.14), w - Inches(3.35), h - Inches(0.28),
        desc, sz=desc_sz, color=DARK_TEXT, wrap=True)

def footer(slide, text):
    txt(slide, Inches(0.5), H - Inches(0.38), W - Inches(1.0), Inches(0.3),
        text, sz=9, color=OFF_WHITE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_COVER)

# White title block on left (clean Snowflake style)
rect(sl, Inches(0.0), Inches(0.0), Inches(5.8), H, fill=CARD_WHITE)
rect(sl, Inches(5.8), Inches(0.0), Inches(0.06), H, fill=SF_BLUE)

# Title area
txt(sl, Inches(0.4), Inches(0.5), Inches(5.1), Inches(0.55),
    "SUPER MARIO", sz=13, color=SF_NAVY, bold=True)
txt(sl, Inches(0.4), Inches(1.0), Inches(5.1), Inches(1.1),
    "SPCS Telemetry", sz=40, color=SF_NAVY, bold=True)
rect(sl, Inches(0.4), Inches(2.18), Inches(1.5), Inches(0.05), fill=SF_ORANGE)
txt(sl, Inches(0.4), Inches(2.3), Inches(5.1), Inches(0.55),
    "Real-time Game Analytics on Snowflake", sz=15, color=SF_NAVY)
txt(sl, Inches(0.4), Inches(2.9), Inches(5.1), Inches(0.4),
    "Snowpark Container Services  ·  OpenTelemetry  ·  Interactive Tables  ·  Cortex AI",
    sz=11, color=DARK_TEXT)

# Accent pills
pill_data = [("🐳 SPCS", SF_ORANGE), ("📡 OTel", SF_TEAL), ("⚡ Tables", SF_PURPLE), ("🤖 Cortex", SF_PINK)]
for i, (lbl, col) in enumerate(pill_data):
    px = Inches(0.4 + i * 1.3)
    rect(sl, px, Inches(3.55), Inches(1.18), Inches(0.38), fill=col, radius=0.1)
    txt(sl, px + Inches(0.1), Inches(3.6), Inches(1.0), Inches(0.28),
        lbl, sz=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(sl, Inches(0.4), Inches(4.15), Inches(5.1), Inches(0.3),
    "Solutions Engineering Demo  ·  20 min including live demo",
    sz=11, color=SF_NAVY, italic=True)
txt(sl, Inches(0.4), Inches(4.55), Inches(5.1), Inches(0.3),
    "github.com/sfc-gh-tdahlberg/mario-spcs-telemetry",
    sz=10, color=SF_BLUE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
slide_title(sl, "Agenda", subtitle="20 minutes · architecture + live demo + learnings")

items = [
    ("01", "What We Built",            "Five components overview · 2 min",              SF_TEAL),
    ("02", "Architecture Deep Dive",   "Full data flow from browser to dashboard · 3 min", SF_ORANGE),
    ("03", "Player Identity",          "How SPCS auth flows into telemetry · 2 min",    SF_PURPLE),
    ("04", "Real-time Pipeline",       "Interactive Tables + Cortex AI · 2 min",        SF_PINK),
    ("🎮", "LIVE DEMO",                "Play · watch · ask Cortex · 10 min",            SF_ORANGE),
    ("05", "Key Learnings",            "What surprised us building this · 1 min",       SF_TEAL),
]

# Two columns
for i, (num, title, sub, col) in enumerate(items):
    row = i % 3; col_idx = i // 3
    lx = Inches(0.45 + col_idx * 5.0)
    ty = Inches(1.12 + row * 1.37)
    rect(sl, lx, ty, Inches(4.6), Inches(1.2), fill=CARD_WHITE, radius=0.05)
    rect(sl, lx, ty, Inches(4.6), Inches(0.06), fill=col, radius=0.05)
    # Number circle
    sh = sl.shapes.add_shape(MSO_SHAPE.OVAL, lx + Inches(0.15), ty + Inches(0.25),
                              Inches(0.5), Inches(0.5))
    sh.fill.solid(); sh.fill.fore_color.rgb = col; sh.line.fill.background()
    sh.text_frame.paragraphs[0].text = num
    sh.text_frame.paragraphs[0].font.size = Pt(14)
    sh.text_frame.paragraphs[0].font.bold = True
    sh.text_frame.paragraphs[0].font.color.rgb = WHITE
    sh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    txt(sl, lx + Inches(0.8), ty + Inches(0.22), Inches(3.65), Inches(0.36),
        title, sz=14, color=SF_NAVY, bold=True)
    txt(sl, lx + Inches(0.8), ty + Inches(0.62), Inches(3.65), Inches(0.45),
        sub, sz=11, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT WE BUILT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
slide_title(sl, "What We Built", subtitle="Five components · one Snowflake account · zero external infrastructure")

components = [
    ("🎮", "Game on SPCS",         ["HTML5 Super Mario in a Docker", "container with public ingress,", "authentication, and auto-resume"],                    SF_ORANGE),
    ("📡", "Telemetry Pipeline",   ["Browser JS → nginx → Python", "sidecar → OpenTelemetry OTLP", "gRPC → SPCS Event Table"],                             SF_TEAL),
    ("⚡", "Interactive Tables",   ["6 live tables at 1-min lag", "Interactive Warehouse (IWH)", "sub-second query latency"],                               SF_PURPLE),
    ("🤖", "Cortex AI",            ["Semantic View + 10 VQRs", "FastGen-generated queries", "Cortex Agent for NL analytics"],                              SF_PINK),
    ("📊", "Dashboards",           ["Streamlit SiS + React Next.js", "Leaderboard, metrics, events,", "animated data pipeline tab"],                       SF_ORANGE),
]

for i, (icon, title, body, col) in enumerate(components):
    x = Inches(0.3 + i * 1.88)
    white_card(sl, x, Inches(1.12), Inches(1.78), Inches(4.0),
               f"{icon}  {title}", body, accent=col, title_sz=12, body_sz=11)

footer(sl, "All data stays in Snowflake · No external APIs · No extra accounts needed")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
slide_title(sl, "Architecture", subtitle="Browser → SPCS Container → Event Table → Interactive Tables → Dashboards")

# Flow nodes as white cards
nodes = [
    (Inches(0.18), Inches(1.2), Inches(1.55), Inches(1.1), "🌐", "Browser", "telemetry.js", SF_TEAL),
    (Inches(2.0), Inches(1.2), Inches(1.55), Inches(1.1), "🔐", "SPCS Ingress", "injects user header", SF_ORANGE),
    (Inches(3.82), Inches(1.2), Inches(1.55), Inches(1.1), "⚙️", "nginx :8080", "routes traffic", SF_BLUE),
    (Inches(3.82), Inches(2.7), Inches(1.55), Inches(1.0), "🐱", "Tomcat :8888", "game server", SF_PINK),
    (Inches(5.68), Inches(1.2), Inches(1.62), Inches(1.1), "🐍", "Python Sidecar", "OTel :9090", SF_PURPLE),
    (Inches(7.6), Inches(1.2), Inches(2.12), Inches(1.1), "📋", "Event Table", "event_db.event_sh", SF_TEAL),
    (Inches(7.6), Inches(2.7), Inches(2.12), Inches(1.0), "⚡", "Interactive Tables", "DIS_MARIO 6×", SF_ORANGE),
    (Inches(7.6), Inches(4.1), Inches(2.12), Inches(0.85), "🤖", "Cortex AI", "Agent + Semantic View", SF_PURPLE),
]

for l, t, w, h, icon, title, sub, col in nodes:
    rect(sl, l, t, w, h, fill=CARD_WHITE, radius=0.05)
    rect(sl, l, t, w, Inches(0.05), fill=col, radius=0.04)
    txt(sl, l + Inches(0.1), t + Inches(0.1), w - Inches(0.2), Inches(0.36),
        f"{icon}  {title}", sz=12, color=SF_NAVY, bold=True)
    txt(sl, l + Inches(0.1), t + Inches(0.48), w - Inches(0.2), Inches(0.35),
        sub, sz=10, color=DARK_TEXT)

# Arrow labels
arrows = [
    (Inches(1.75), Inches(1.65), "GET pixel"),
    (Inches(3.58), Inches(1.65), "header"),
    (Inches(5.46), Inches(1.65), "OTLP gRPC"),
    (Inches(7.6),  Inches(2.35), "1-min lag"),
    (Inches(7.6),  Inches(3.75), "SQL views"),
    (Inches(4.62), Inches(2.22), "proxy /"),
]
for ax, ay, lbl in arrows:
    txt(sl, ax, ay - Inches(0.18), Inches(0.2), Inches(0.2), "▶", sz=10, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, ax - Inches(0.1), ay + Inches(0.04), Inches(0.42), Inches(0.2),
        lbl, sz=7, color=OFF_WHITE, italic=True, align=PP_ALIGN.CENTER)

# Consumers footer row
for i, (icon, lbl, col) in enumerate([("📊","Streamlit",SF_TEAL),("⚛️","React App",SF_ORANGE),("🤖","Cortex Agent",SF_PURPLE)]):
    cx = Inches(7.62 + i * 0.0)
    pass  # already represented in nodes above

footer(sl, "Single SPCS container: nginx → Tomcat (game) + Python sidecar (OTel) · Compute Pool: CPU_X64_XS")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CONTAINER STACK
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
slide_title(sl, "Container Stack", subtitle="One Docker image · three processes · everything Snowflake needs")

layers = [
    (SF_TEAL,   "🌐  nginx  ·  Port 8080",
     "Public SPCS ingress endpoint. Exact location = /telemetry (avoids catching /telemetry.js). Forwards Sf-Context-Current-User header to sidecar."),
    (SF_ORANGE, "🐱  Tomcat 9  ·  Port 8888",
     "HTML5 Super Mario game. telemetry.js injected at Docker build time. jQuery bundled locally — SPCS egress blocks external CDN URLs."),
    (SF_PURPLE, "🐍  Python Sidecar  ·  Port 9090",
     "OpenTelemetry TracerProvider + MeterProvider. GET tracking pixel /telemetry?d=... · /whoami endpoint. BatchSpanProcessor exports every 5 seconds."),
    (SF_PINK,   "📡  OpenTelemetry OTLP gRPC",
     "Spans → record_attributes · Metrics → 90+ platform types · Logs → service messages. All land in event_db.event_sh.my_events automatically."),
]

for i, (col, title, desc) in enumerate(layers):
    y = Inches(1.18 + i * 1.06)
    rect(sl, Inches(0.4), y, W - Inches(0.8), Inches(0.92), fill=CARD_WHITE, radius=0.04)
    rect(sl, Inches(0.4), y, Inches(0.08), Inches(0.92), fill=col)
    txt(sl, Inches(0.65), y + Inches(0.1), Inches(3.5), Inches(0.38),
        title, sz=15, color=SF_NAVY, bold=True)
    txt(sl, Inches(0.65), y + Inches(0.52), W - Inches(1.2), Inches(0.36),
        desc, sz=11, color=DARK_TEXT, wrap=True)

footer(sl, "Resources: 500m–1 vCPU · 512Mi–2Gi RAM · AUTO_RESUME = TRUE · platformMonitor: system, network, storage")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PLAYER IDENTITY
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
slide_title(sl, "Player Identity Flow", subtitle="From SPCS login to named leaderboard entry")

# 7-step flow
steps = [
    ("1", "Ingress Auth",    "Cookie",                SF_TEAL),
    ("2", "Inject Header",   "Sf-Context-\nCurrent-User", SF_ORANGE),
    ("3", "nginx Forward",   "proxy_set_\nheader",    SF_PURPLE),
    ("4", "/whoami",         "_get_player_\nname()",  SF_PINK),
    ("5", "Browser Fetch",   "XHR on\ngame init",     SF_TEAL),
    ("6", "Event Payload",   "player_name\nin send()", SF_ORANGE),
    ("7", "OTel Span",       "record_\nattributes",   SF_PURPLE),
]

step_w = Inches(1.27)
for i, (num, title, sub, col) in enumerate(steps):
    x = Inches(0.28 + i * 1.36)
    rect(sl, x, Inches(1.1), step_w, Inches(1.6), fill=CARD_WHITE, radius=0.05)
    rect(sl, x, Inches(1.1), step_w, Inches(0.055), fill=col, radius=0.04)
    sh = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.39), Inches(1.2),
                             Inches(0.48), Inches(0.48))
    sh.fill.solid(); sh.fill.fore_color.rgb = col; sh.line.fill.background()
    sh.text_frame.paragraphs[0].text = num
    sh.text_frame.paragraphs[0].font.size = Pt(16)
    sh.text_frame.paragraphs[0].font.bold = True
    sh.text_frame.paragraphs[0].font.color.rgb = WHITE
    sh.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    txt(sl, x + Inches(0.08), Inches(1.78), step_w - Inches(0.16), Inches(0.36),
        title, sz=12, color=SF_NAVY, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.08), Inches(2.18), step_w - Inches(0.16), Inches(0.44),
        sub, sz=10, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    if i < 6:
        txt(sl, x + step_w, Inches(1.82), Inches(0.1), Inches(0.3),
            "▶", sz=12, color=WHITE, align=PP_ALIGN.CENTER)

# Priority table
ty = Inches(2.95)
rect(sl, Inches(0.4), ty, W - Inches(0.8), Inches(1.6), fill=CARD_WHITE, radius=0.05)
rect(sl, Inches(0.4), ty, W - Inches(0.8), Inches(0.055), fill=SF_NAVY, radius=0.04)
txt(sl, Inches(0.6), ty + Inches(0.1), Inches(4), Inches(0.36),
    "Priority Chain", sz=13, color=SF_NAVY, bold=True)
chain = [
    ("1", "Sf-Context-Current-User header",    "Injected by SPCS ingress — most reliable",  SF_TEAL),
    ("2", "player_name in event payload",       "Sent by browser after /whoami fetch — reliable fallback", SF_ORANGE),
    ("3", "\"unknown\"",                        "Last resort if both header and browser sources are empty", SF_PINK),
]
for j, (pri, src, note, col) in enumerate(chain):
    ry = ty + Inches(0.55 + j * 0.33)
    rect(sl, Inches(0.5), ry, Inches(0.28), Inches(0.26), fill=col, radius=0.05)
    txt(sl, Inches(0.52), ry + Inches(0.04), Inches(0.24), Inches(0.2),
        pri, sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, Inches(0.88), ry + Inches(0.04), Inches(2.8), Inches(0.22), src, sz=11, color=SF_NAVY, bold=True)
    txt(sl, Inches(3.75), ry + Inches(0.04), Inches(5.8), Inches(0.22), note, sz=10, color=DARK_TEXT)

footer(sl, "Root cause: telemetry.js fetched /whoami for the banner but never included player_name in event payloads — fixed by adding it to every send()")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — REAL-TIME PIPELINE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
slide_title(sl, "Real-time Data Pipeline", subtitle="DIS_MARIO: Interactive Tables + IWH + Semantic View + Cortex Agent")

tables = [
    ("GAME_EVENTS_LIVE",      "EVENT_TIME",    "All game spans + player_name"),
    ("EVENT_TIMELINE_LIVE",   "MINUTE",        "Event counts per minute"),
    ("KEY_PRESSES_LIVE",      "KEY_NAME",      "Key press aggregates"),
    ("DEATHS_BY_LEVEL_LIVE",  "LEVEL",         "Deaths per level"),
    ("POWERUPS_LIVE",         "POWERUP_TYPE",  "Powerup spawn counts"),
    ("PLAYER_SESSIONS_LIVE",  "SESSION_START", "Per-player session summaries"),
]

for i, (name, cluster, desc) in enumerate(tables):
    y = Inches(1.1 + i * 0.72)
    rect(sl, Inches(0.38), y, Inches(5.6), Inches(0.64), fill=CARD_WHITE, radius=0.04)
    rect(sl, Inches(0.38), y, Inches(0.07), Inches(0.64), fill=SF_PURPLE)
    txt(sl, Inches(0.58), y + Inches(0.08), Inches(2.7), Inches(0.28),
        name, sz=12, color=SF_NAVY, bold=True)
    txt(sl, Inches(0.58), y + Inches(0.38), Inches(2.7), Inches(0.22),
        desc, sz=10, color=DARK_TEXT)
    rect(sl, Inches(3.4), y + Inches(0.15), Inches(1.1), Inches(0.3),
         fill=SF_PURPLE, radius=0.06)
    txt(sl, Inches(3.4), y + Inches(0.15), Inches(1.1), Inches(0.3),
        f"⚡ {cluster}", sz=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, Inches(4.6), y + Inches(0.2), Inches(1.3), Inches(0.2),
        "lag: 1 min", sz=9, color=SF_TEAL, bold=True)

# Right panel
right_items = [
    (Inches(6.2), Inches(1.1), Inches(3.5), Inches(1.32), "⚡ Interactive Warehouse",
     ["DIS_MARIO_IWH · XSMALL", "24h auto-suspend · sub-second latency", "For all 6 tables (ALTER INTERACTIVE TABLE … SET WAREHOUSE)"], SF_ORANGE),
    (Inches(6.2), Inches(2.52), Inches(3.5), Inches(1.32), "🧠 Semantic View",
     ["DIS_MARIO.PUBLIC.MARIO_TELEMETRY", "6 tables · 2 joins · 10 VQRs", "Created via SYSTEM$CORTEX_ANALYST_FAST_GENERATION"], SF_TEAL),
    (Inches(6.2), Inches(3.94), Inches(3.5), Inches(1.1), "🤖 Cortex Agent",
     ["DIS_MARIO.PUBLIC.MARIO_INTELLIGENCE", "model: auto · 900s / 400k tokens", "NL → SQL via MARIO_TELEMETRY"], SF_PURPLE),
]
for lx, ty, w, h, title, body, col in right_items:
    white_card(sl, lx, ty, w, h, title, body, accent=col, title_sz=12, body_sz=10)

footer(sl, "CREATE INTERACTIVE TABLE (not DYNAMIC) · ALTER INTERACTIVE TABLE … SET WAREHOUSE = DIS_MARIO_IWH")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — DASHBOARDS & AI
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
slide_title(sl, "Dashboards & AI", subtitle="Four ways to consume the real-time telemetry data")

dashboards = [
    ("📊", "Streamlit — Legacy",
     "MARIO_DB.PUBLIC.MARIO_TELEMETRY_DASHBOARD",
     ["Event table views", "Leaderboard: last 24h only", "5 tabs: Leaderboard, Events,", "Deaths, Controls, Platform Metrics"],
     SF_TEAL, Inches(0.3), Inches(1.1)),
    ("⚡", "Streamlit — Live",
     "DIS_MARIO.PUBLIC.DIS_MARIO_TELEMETRY_DASHBOARD",
     ["Interactive Tables backend", "10s cache TTL · near-realtime", "Player dropdown filter", "6 KPI cards + Snowflake branding"],
     SF_ORANGE, Inches(0.3), Inches(3.18)),
    ("⚛️", "React App",
     "Next.js 16.2.3 · port 3456",
     ["JWT auth locally, OAuth in SPCS", "5s polling · Recharts charts", "Animated Data Pipeline tab", "Snowflake logo + Cortex Code badge"],
     SF_PURPLE, Inches(5.1), Inches(1.1)),
    ("🤖", "Cortex Agent",
     "DIS_MARIO.PUBLIC.MARIO_INTELLIGENCE",
     ["NL → SQL via MARIO_TELEMETRY", "model: auto · 900s budget", "Try: 'Who died most on 1-1?'", "'What's the avg session time?'"],
     SF_PINK, Inches(5.1), Inches(3.18)),
]

for icon, title, subtitle, body, col, lx, ty in dashboards:
    w, h = Inches(4.55), Inches(1.92)
    rect(sl, lx, ty, w, h, fill=CARD_WHITE, radius=0.05)
    rect(sl, lx, ty, w, Inches(0.06), fill=col, radius=0.04)
    txt(sl, lx + Inches(0.14), ty + Inches(0.12), Inches(0.42), Inches(0.42), icon, sz=22)
    txt(sl, lx + Inches(0.65), ty + Inches(0.1), w - Inches(0.8), Inches(0.32),
        title, sz=13, color=SF_NAVY, bold=True)
    txt(sl, lx + Inches(0.65), ty + Inches(0.44), w - Inches(0.8), Inches(0.24),
        subtitle, sz=9, color=SF_NAVY, italic=True)
    multiline_txt(sl, lx + Inches(0.14), ty + Inches(0.76), w - Inches(0.28),
                  h - Inches(0.82), body, sz=10, color=DARK_TEXT, spacing=2)

footer(sl, "All dashboards: Snowflake logotype SVG · Powered by Cortex Code badge · Animated polar bear 🐻‍❄️")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SNOWFLAKE FEATURES
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
slide_title(sl, "Snowflake Features Demonstrated",
            subtitle="Every component in this demo is a production Snowflake feature")

features = [
    (SF_TEAL,   "Snowpark Container Services",
     "Docker container · public ingress · platform monitoring · Sf-Context-Current-User header injection"),
    (SF_ORANGE, "SPCS Event Tables",
     "Native OTel store for spans, metrics, logs · OTLP gRPC exporter · 90+ built-in platform metric types"),
    (SF_PURPLE, "Interactive Tables",
     "Sub-second query latency · 1-min target lag · cluster keys · Interactive Warehouse · near-realtime dashboards"),
    (SF_PINK,   "Cortex Analyst",
     "Semantic View over 6 tables · FastGen VQR generation · SYSTEM$CREATE_SEMANTIC_VIEW_FROM_YAML"),
    (SF_TEAL,   "Cortex Agents",
     "Custom orchestration + response instructions · MARIO_INTELLIGENCE agent · cortex_analyst_text_to_sql tool"),
    (SF_ORANGE, "Streamlit in Snowsight",
     "Container runtime · SYSTEM_COMPUTE_POOL_CPU · pip packages · snow CLI deploy · player filter"),
    (SF_PURPLE, "Semi-structured JSON",
     "record:name::STRING · record_attributes:player_name::STRING · VARIANT extraction at query time"),
]

for i, (col, title, desc) in enumerate(features):
    y = Inches(1.1 + i * 0.63)
    rect(sl, Inches(0.38), y, W - Inches(0.76), Inches(0.56), fill=CARD_WHITE, radius=0.04)
    rect(sl, Inches(0.38), y, Inches(0.07), Inches(0.56), fill=col)
    txt(sl, Inches(0.58), y + Inches(0.08), Inches(3.0), Inches(0.28), title, sz=13, color=SF_NAVY, bold=True)
    txt(sl, Inches(0.58), y + Inches(0.32), W - Inches(1.1), Inches(0.22), desc, sz=10, color=DARK_TEXT)

footer(sl, "Zero external infrastructure · All data stays in Snowflake · Everything demos in a single account")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DEMO DIVIDER
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_DIVIDER)  # cloud sky background — perfect for Mario!

# White content block
rect(sl, Inches(0.0), Inches(0.0), Inches(5.5), H, fill=CARD_WHITE)
rect(sl, Inches(5.5), Inches(0.0), Inches(0.06), H, fill=SF_ORANGE)

txt(sl, Inches(0.4), Inches(0.45), Inches(4.8), Inches(0.55),
    "🎮  LIVE DEMO", sz=36, color=SF_NAVY, bold=True)
rect(sl, Inches(0.4), Inches(1.12), Inches(2.2), Inches(0.05), fill=SF_ORANGE)
txt(sl, Inches(0.4), Inches(1.25), Inches(4.8), Inches(0.38),
    "10 minutes · play · watch · ask Cortex", sz=14, color=SF_NAVY, italic=True)

demo_steps = [
    (SF_TEAL,   "Step 1", "Open game · press S to start · play Mario"),
    (SF_ORANGE, "Step 2", "Watch events → sidecar → event table in logs"),
    (SF_PURPLE, "Step 3", "Refresh Streamlit · see name in leaderboard"),
    (SF_PINK,   "Step 4", "Ask Cortex Agent: 'Who died most on 1-1?'"),
]
for i, (col, label, action) in enumerate(demo_steps):
    sy = Inches(1.82 + i * 0.86)
    rect(sl, Inches(0.4), sy, Inches(4.8), Inches(0.72), fill=col, radius=0.05)
    txt(sl, Inches(0.6), sy + Inches(0.1), Inches(0.9), Inches(0.3),
        label, sz=11, color=WHITE, bold=True)
    txt(sl, Inches(1.6), sy + Inches(0.1), Inches(3.5), Inches(0.5),
        action, sz=12, color=WHITE)

txt(sl, Inches(0.4), H - Inches(0.45), Inches(4.8), Inches(0.32),
    "ei53mb-sfseeurope-eu-demo200.snowflakecomputing.app",
    sz=9, color=SF_BLUE, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — KEY LEARNINGS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_CONTENT)
slide_title(sl, "Key Learnings", subtitle="Five things that surprised us building this")

learnings = [
    (SF_TEAL,   "🐳  Always use --no-cache",
     ["Docker silently reuses stale cached layers.", "suspend/resume does NOT re-pull the image.", "Use ALTER SERVICE FROM SPECIFICATION."]),
    (SF_ORANGE, "📍  nginx exact match",
     ["location /telemetry is a prefix match —", "it catches /telemetry.js too.", "Always use: location = /telemetry"]),
    (SF_PURPLE, "👤  Browser must send player_name",
     ["Sidecar had header extraction but browser", "was calling /whoami for banner only.", "Fix: add player_name to every send() call."]),
    (SF_PINK,   "📦  Bundle all dependencies",
     ["SPCS blocks external egress.", "CDN URLs fail silently at runtime.", "Bundle jQuery inside the Docker image."]),
    (SF_TEAL,   "🔑  JWT for headless auth",
     ["EXTERNALBROWSER fails in Node.js.", "Use SNOWFLAKE_JWT keypair auth.", "snow spcs image-registry login for MFA."]),
]

for i, (col, title, body) in enumerate(learnings):
    row = i // 3; col_idx = i % 3
    if i < 3:
        lx = Inches(0.3 + col_idx * 3.22)
        ty = Inches(1.1)
    else:
        lx = Inches(0.3 + (i-3) * 3.22)
        ty = Inches(3.28) if i < 5 else Inches(3.28)
        if i == 4: lx = Inches(3.52)

    white_card(sl, lx, ty, Inches(3.08), Inches(2.0), title, body,
               accent=col, title_sz=12, body_sz=11)

footer(sl, "Full details in README.md and docs/index.html · github.com/sfc-gh-tdahlberg/mario-spcs-telemetry")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(LAY_THANKS)

# White left panel on the thank you slide
rect(sl, Inches(0.0), Inches(0.0), Inches(5.8), H, fill=CARD_WHITE)
rect(sl, Inches(5.8), Inches(0.0), Inches(0.06), H, fill=SF_BLUE)

txt(sl, Inches(0.45), Inches(0.55), Inches(5.1), Inches(0.7),
    "Thanks for playing! 🍄", sz=32, color=SF_NAVY, bold=True)
rect(sl, Inches(0.45), Inches(1.4), Inches(2.0), Inches(0.05), fill=SF_TEAL)
txt(sl, Inches(0.45), Inches(1.55), Inches(5.1), Inches(0.4),
    "Questions welcome — let's talk SPCS, Interactive Tables, Cortex AI",
    sz=13, color=DARK_TEXT, italic=True)

links = [
    ("🎮", "Game URL",       "ei53mb-sfseeurope-eu-demo200\n.snowflakecomputing.app", SF_TEAL),
    ("💻", "GitHub",         "sfc-gh-tdahlberg/\nmario-spcs-telemetry",              SF_ORANGE),
    ("📖", "HTML Docs",      "docs/index.html\nfull architecture reference",         SF_PURPLE),
    ("🤖", "Cortex Agent",   "DIS_MARIO.PUBLIC\n.MARIO_INTELLIGENCE",                SF_PINK),
]

for i, (icon, label, val, col) in enumerate(links):
    row, c = divmod(i, 2)
    lx = Inches(0.45 + c * 2.7)
    ty = Inches(2.1 + row * 1.32)
    rect(sl, lx, ty, Inches(2.5), Inches(1.18), fill=col, radius=0.06)
    txt(sl, lx + Inches(0.15), ty + Inches(0.1), Inches(0.45), Inches(0.42), icon, sz=20)
    txt(sl, lx + Inches(0.65), ty + Inches(0.1), Inches(1.72), Inches(0.3),
        label, sz=12, color=WHITE, bold=True)
    multiline_txt(sl, lx + Inches(0.15), ty + Inches(0.55), Inches(2.22),
                  Inches(0.6), val.split("\n"), sz=9, color=WHITE, spacing=1)

txt(sl, Inches(0.45), H - Inches(0.42), Inches(5.1), Inches(0.28),
    "Built with ✨ Cortex Code · Snowflake Solutions Engineering",
    sz=10, color=SF_NAVY, italic=True)


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"✅  Saved → {OUTPUT}  ({len(prs.slides)} slides)")
